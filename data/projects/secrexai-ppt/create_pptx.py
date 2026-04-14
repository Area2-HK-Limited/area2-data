#!/usr/bin/env python3
"""
SecrexAI Introduction PPTX Generator
For Area2 Sales - Traditional Chinese (HK)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Color scheme - Modern Blue/Teal
PRIMARY_BLUE = RGBColor(0, 102, 153)      # #006699
TEAL = RGBColor(0, 153, 153)               # #009999
LIGHT_BLUE = RGBColor(204, 229, 255)      # #CCE5FF
DARK_TEXT = RGBColor(51, 51, 51)          # #333333
WHITE = RGBColor(255, 255, 255)
ACCENT_ORANGE = RGBColor(255, 153, 0)     # #FF9900

# Image paths
IMG_DIR = "/home/ubuntu/.openclaw/workspace/data/projects/secrexai-ppt/images"
OUTPUT_PATH = "/home/ubuntu/.openclaw/workspace/data/projects/secrexai-ppt/SecrexAI_Introduction.pptx"

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background image
    img_path = f"{IMG_DIR}/slide1_hero.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(13.33), height=Inches(7.5))
    
    # Dark overlay for text readability
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5), Inches(13.33), Inches(2.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 51, 76)
    overlay.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SecrexAI 介紹"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "香港 SME 專屬 WhatsApp AI 助手"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 204, 204)
    p.alignment = PP_ALIGN.CENTER
    
    # Company tag
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.5))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Area2 呈獻"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER


def add_what_is_slide(prs):
    """Slide 2: 什麼是 SecrexAI"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "什麼是 SecrexAI?"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content - 4 feature boxes in 2x2 grid
    features = [
        ("24小時 AI 助手", "住喺你 WhatsApp"),
        ("廣東話原生", "零學習成本"),
        ("永久記憶", "越用越就手"),
        ("60+ 技能", "一个平台解決所有需求"),
    ]
    
    positions = [(0.5, 1.8), (6.9, 1.8), (0.5, 4.2), (6.9, 4.2)]
    colors = [PRIMARY_BLUE, TEAL, TEAL, PRIMARY_BLUE]
    
    for i, (title, desc) in enumerate(features):
        x, y = positions[i]
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.9), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()
        
        # Title text
        t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(5.3), Inches(0.8))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Desc text
        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.1), Inches(5.3), Inches(0.7))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(230, 230, 230)


def add_comparison_slide(prs):
    """Slide 3: Comparison with General AI"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "同一般 AI 助手比較"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Comparison image
    img_path = f"{IMG_DIR}/slide3_comparison.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12.3), height=Inches(5.5))
    
    # Table overlay
    table_data = [
        ("一般 AI", "SecrexAI"),
        ("只答你", "幫你做"),
        ("獨立對話，無記憶", "永久記憶"),
        ("要你自己郁手做", "自動幫你完成"),
        ("無提醒", "24/7 自動跟進"),
    ]
    
    # Left column - General AI (red-ish)
    left_x = 0.5
    for i, (left, right) in enumerate(table_data):
        y = 1.6 + (i * 1.0)
        if i == 0:
            # Header
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_x), Inches(y), Inches(5.8), Inches(0.8))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(180, 60, 60)
            box.line.fill.background()
            
            t_box = slide.shapes.add_textbox(Inches(left_x + 0.2), Inches(y + 0.15), Inches(5.4), Inches(0.5))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = left
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        else:
            t_box = slide.shapes.add_textbox(Inches(left_x + 0.2), Inches(y + 0.1), Inches(5.4), Inches(0.7))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"❌ {left}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(180, 60, 60)
    
    # Right column - SecrexAI (green-ish)
    right_x = 6.9
    for i, (left, right) in enumerate(table_data):
        y = 1.6 + (i * 1.0)
        if i == 0:
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x), Inches(y), Inches(5.8), Inches(0.8))
            box.fill.solid()
            box.fill.fore_color.rgb = TEAL
            box.line.fill.background()
            
            t_box = slide.shapes.add_textbox(Inches(right_x + 0.2), Inches(y + 0.15), Inches(5.4), Inches(0.5))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = right
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        else:
            t_box = slide.shapes.add_textbox(Inches(right_x + 0.2), Inches(y + 0.1), Inches(5.4), Inches(0.7))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✅ {right}"
            p.font.size = Pt(22)
            p.font.color.rgb = TEAL


def add_features_slide(prs):
    """Slide 4: Core Features"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心功能"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Features list
    features = [
        ("60+ 技能", "搜尋、問答、提醒、整理文件"),
        ("永久記憶系統", "越用越了解你"),
        ("排程提醒功能", "自動跟進重要事項"),
        ("WhatsApp 群組整合", "群組助手功能"),
        ("自動生成內容", "文案、圖片、分析報告"),
    ]
    
    for i, (title, desc) in enumerate(features):
        y = 1.6 + (i * 1.05)
        
        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.1), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY_BLUE
        circle.line.fill.background()
        
        # Number
        num_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(0.7), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(1.6), Inches(y + 0.05), Inches(5), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        
        # Description
        d_box = slide.shapes.add_textbox(Inches(1.6), Inches(y + 0.5), Inches(10), Inches(0.5))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(100, 100, 100)


def add_pricing_slide(prs):
    """Slide 5: Pricing"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background image
    img_path = f"{IMG_DIR}/slide5_pricing.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(13.33), height=Inches(7.5))
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "定價方案"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Pricing cards
    plans = [
        ("Starter 個人版", "HK$298/月", False),
        ("Business 商務版", "HK$688/月", True),  # Most popular
        ("Enterprise 企業版", "度身訂造", False),
    ]
    
    for i, (name, price, popular) in enumerate(plans):
        x = 0.8 + (i * 4.2)
        y = 2.0
        
        # Card background
        if popular:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(4.5))
            card.fill.solid()
            card.fill.fore_color.rgb = PRIMARY_BLUE
            card.line.fill.background()
        else:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(4.5))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(240, 240, 240)
            card.line.color.rgb = RGBColor(200, 200, 200)
        
        # Popular badge
        if popular:
            badge = slide.shapes.add_textbox(Inches(x), Inches(y - 0.3), Inches(3.8), Inches(0.5))
            tf = badge.text_frame
            p = tf.paragraphs[0]
            p.text = "⭐ 最受歡迎"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ORANGE
            p.alignment = PP_ALIGN.CENTER
        
        # Plan name
        name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.5), Inches(3.4), Inches(1))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE if not popular else WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Price
        price_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.5), Inches(3.4), Inches(1))
        tf = price_box.text_frame
        p = tf.paragraphs[0]
        p.text = price
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEAL if not popular else WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Free trial note
    trial_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
    tf = trial_box.text_frame
    p = tf.paragraphs[0]
    p.text = "所有方案：14日免費試用"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Button shape
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(6.7), Inches(4), Inches(0.6))
    btn.fill.solid()
    btn.fill.fore_color.rgb = ACCENT_ORANGE
    btn.line.fill.background()


def add_target_audience_slide(prs):
    """Slide 6: Target Audience"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "適合邊類人?"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Audience items
    audiences = [
        ("中小企老闆", "用 AI 提升效率"),
        ("零售店 / 餐廳東主", "自動回覆客人查詢"),
        ("美容院 / 服務行業", "智能預約管理"),
        ("地產 / 會計 / 物流", "專業客戶跟進"),
    ]
    
    for i, (title, desc) in enumerate(audiences):
        x = 0.5 + (i * 3.2)
        y = 2.0
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 0 else RGBColor(230, 245, 245)
        card.line.fill.background()
        
        # Icon placeholder (circle)
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.95), Inches(y + 0.5), Inches(1), Inches(1))
        icon.fill.solid()
        icon.fill.fore_color.rgb = PRIMARY_BLUE if i % 2 == 0 else TEAL
        icon.line.fill.background()
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.8), Inches(2.7), Inches(1))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # Desc
        d_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 2.8), Inches(2.7), Inches(1.5))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER


def add_security_slide(prs):
    """Slide 7: Security & Trust"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background image
    img_path = f"{IMG_DIR}/slide7_security.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(13.33), height=Inches(7.5))
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "為什麼選擇 SecrexAI?"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Trust items
    items = [
        ("🏆 ISO 27001 / SOC 2", "企業級資安認證"),
        ("🇭🇰 PDPO 合規", "香港公司運營"),
        ("☁️ 數據存亞太區雲端", "安全可靠"),
        ("⚡ 30分鐘個人化 Onboarding", "快速上手"),
        ("💬 WhatsApp 廣東話技術支援", "貼心服務"),
    ]
    
    for i, (title, desc) in enumerate(items):
        y = 1.5 + (i * 1.0)
        
        # Icon
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.1), Inches(0.7), Inches(0.7))
        icon.fill.solid()
        icon.fill.fore_color.rgb = TEAL
        icon.line.fill.background()
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.05), Inches(10), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        
        # Desc
        d_box = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.5), Inches(10), Inches(0.4))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(100, 100, 100)


def add_cta_slide(prs):
    """Slide 8: Call to Action"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Gradient background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BLUE
    bg.line.fill.background()
    
    # Decorative circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    
    # Main CTA text
    cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
    tf = cta_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "立即開始"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Highlight text
    highlight_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1))
    tf = highlight_box.text_frame
    p = tf.paragraphs[0]
    p.text = "14 日免費試用"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(1.5))
    tf = contact_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "聯絡 Area2 獲取更多資訊"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "WhatsApp: 9565 1459 | info@area2.com"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(180, 200, 240)
    p2.alignment = PP_ALIGN.CENTER
    
    # Company
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Area2 Limited"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(150, 180, 220)
    p.alignment = PP_ALIGN.CENTER


def main():
    print("🚀 Creating SecrexAI Introduction PPTX...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # Create all slides
    print("📝 Creating Slide 1: Title...")
    add_title_slide(prs)
    
    print("📝 Creating Slide 2: What is SecrexAI...")
    add_what_is_slide(prs)
    
    print("📝 Creating Slide 3: Comparison...")
    add_comparison_slide(prs)
    
    print("📝 Creating Slide 4: Core Features...")
    add_features_slide(prs)
    
    print("📝 Creating Slide 5: Pricing...")
    add_pricing_slide(prs)
    
    print("📝 Creating Slide 6: Target Audience...")
    add_target_audience_slide(prs)
    
    print("📝 Creating Slide 7: Security...")
    add_security_slide(prs)
    
    print("📝 Creating Slide 8: Call to Action...")
    add_cta_slide(prs)
    
    # Save
    print(f"💾 Saving to {OUTPUT_PATH}...")
    prs.save(OUTPUT_PATH)
    print(f"✅ Done! Saved to {OUTPUT_PATH}")
    
    # Verify
    if os.path.exists(OUTPUT_PATH):
        size = os.path.getsize(OUTPUT_PATH)
        print(f"📁 File size: {size:,} bytes ({size/1024/1024:.2f} MB)")


if __name__ == "__main__":
    main()
