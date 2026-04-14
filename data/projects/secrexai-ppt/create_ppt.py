#!/usr/bin/env python3
"""Generate SecrexAI Introduction PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_BLUE = RGBColor(0, 102, 153)    # #006699
SECONDARY_TEAL = RGBColor(0, 153, 153)  # #009999
ACCENT_ORANGE = RGBColor(255, 102, 0)    # #FF6600
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Add a text box to slide"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape

def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rectangle(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    """Add bulleted list"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return shape

# ============== SLIDE 1: Title ==============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Background gradient effect (using rectangle)
add_rectangle(slide1, 0, 0, 13.333, 7.5, PRIMARY_BLUE)

# Decorative shapes
add_rectangle(slide1, 0, 5.5, 13.333, 2, SECONDARY_TEAL)

# Title
add_text_box(slide1, 1, 2, 11, 1.5, "SecrexAI 介紹", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, 1, 3.5, 11, 0.8, "香港 SME 專屬 WhatsApp AI 助手", font_size=28, color=WHITE, align=PP_ALIGN.CENTER)

# Area2 logo text
add_text_box(slide1, 0.5, 6.8, 3, 0.5, "Area2 呈獻", font_size=16, color=WHITE, align=PP_ALIGN.LEFT)

# ============== SLIDE 2: What is SecrexAI ==============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_rectangle(slide2, 0, 0, 13.333, 1.2, PRIMARY_BLUE)
add_text_box(slide2, 0.5, 0.3, 12, 0.7, "什麼是 SecrexAI?", font_size=36, bold=True, color=WHITE)

# Content cards
card_y = 1.8
card_height = 1.2
card_gap = 0.15

features = [
    ("24小時 AI 助手", "住喺你 WhatsApp，随时随地帮你做嘢"),
    ("廣東話原生", "零學習成本，直接廣東話對話"),
    ("永久記憶", "越用越就手，記住你嘅工作習慣"),
    ("60+ 技能", "一个平台解決所有需求"),
]

for i, (title, desc) in enumerate(features):
    y = card_y + i * (card_height + card_gap)
    # Card background
    add_rounded_rectangle(slide2, 0.8, y, 11.5, card_height, LIGHT_GRAY)
    # Icon circle
    add_rounded_rectangle(slide2, 1.0, y + 0.25, 0.7, 0.7, SECONDARY_TEAL)
    # Title
    add_text_box(slide2, 2.0, y + 0.15, 9, 0.5, title, font_size=22, bold=True, color=PRIMARY_BLUE)
    # Description
    add_text_box(slide2, 2.0, y + 0.6, 9, 0.5, desc, font_size=16, color=DARK_GRAY)

# Area2 footer
add_text_box(slide2, 0.3, 7.0, 2, 0.3, "Area2", font_size=12, color=SECONDARY_TEAL)

# ============== SLIDE 3: Comparison ==============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_rectangle(slide3, 0, 0, 13.333, 1.2, PRIMARY_BLUE)
add_text_box(slide3, 0.5, 0.3, 12, 0.7, "同一般 AI 助手比較", font_size=36, bold=True, color=WHITE)

# Left column - General AI
add_rounded_rectangle(slide3, 0.8, 1.6, 5.5, 5, RGBColor(255, 230, 230))
add_text_box(slide3, 0.8, 1.8, 5.5, 0.6, "一般 AI 助理", font_size=24, bold=True, color=RGBColor(180, 50, 50), align=PP_ALIGN.CENTER)

general_ai_items = [
    "✗ 你問佢點做，佢淨係答你",
    "✗ 要你自己複製貼上郁手做",
    "✗ 每次對話都係獨立，無記憶",
    "✗ 你唔記得跟進，佢唔會主動提醒",
    "✗ 無替你預訂、寫電郵、查資料",
]
add_bullet_list(slide3, 1.2, 2.6, 5, 3.5, general_ai_items, font_size=17, color=RGBColor(150, 50, 50))

# Right column - SecrexAI
add_rounded_rectangle(slide3, 7.0, 1.6, 5.5, 5, RGBColor(230, 255, 240))
add_text_box(slide3, 7.0, 1.8, 5.5, 0.6, "SecrexAI", font_size=24, bold=True, color=SECONDARY_TEAL, align=PP_ALIGN.CENTER)

secrex_items = [
    "✓ 你叫佢做，佢真係郁手幫你完成",
    "✓ 自動預訂、寫電郵、整理文件",
    "✓ 永久記憶，越用越就手",
    "✓ 佢有電腦喺手，隨時幫你做嘢",
    "✓ 24/7 standby，有手機就得",
]
add_bullet_list(slide3, 7.4, 2.6, 5, 3.5, secrex_items, font_size=17, color=SECONDARY_TEAL)

# Area2 footer
add_text_box(slide3, 0.3, 7.0, 2, 0.3, "Area2", font_size=12, color=SECONDARY_TEAL)

# ============== SLIDE 4: Core Features ==============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_rectangle(slide4, 0, 0, 13.333, 1.2, PRIMARY_BLUE)
add_text_box(slide4, 0.5, 0.3, 12, 0.7, "核心功能", font_size=36, bold=True, color=WHITE)

# Feature grid (2x3)
features4 = [
    ("60+ 技能", "搜尋、問答、提醒、整理文件"),
    ("永久記憶系統", "記住你係誰，越用越聰明"),
    ("排程提醒功能", "自動跟進，永不忘記"),
    ("WhatsApp 群組", "加入群組，集體使用"),
    ("自動生成文案", "一秒生成專業內容"),
    ("圖片分析", "智能分析，即時回應"),
]

for i, (title, desc) in enumerate(features4):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 1.6 + row * 2.6
    
    add_rounded_rectangle(slide4, x, y, 3.8, 2.2, LIGHT_GRAY)
    add_text_box(slide4, x + 0.2, y + 0.3, 3.4, 0.5, title, font_size=20, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide4, x + 0.2, y + 1.0, 3.4, 1, desc, font_size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Area2 footer
add_text_box(slide4, 0.3, 7.0, 2, 0.3, "Area2", font_size=12, color=SECONDARY_TEAL)

# ============== SLIDE 5: Pricing ==============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_rectangle(slide5, 0, 0, 13.333, 1.2, PRIMARY_BLUE)
add_text_box(slide5, 0.5, 0.3, 12, 0.7, "定價方案", font_size=36, bold=True, color=WHITE)

# Three pricing cards
plans = [
    ("Starter", "個人版", "HK$298", "/月", "基本用量\n60+ 技能完整使用\n廣東話原生對話\n永久記憶系統", False),
    ("Business", "商務版", "HK$688", "/月", "Starter 所有功能\n5倍商業用量\nWhatsApp 群組\n優先技術支援", True),
    ("Enterprise", "企業版", "度身訂造", "", "Business 所有功能\n無限成員使用\n自訂 AI 個性\n專屬客戶經理", False),
]

card_width = 3.8
card_gap = 0.5
start_x = (13.333 - 3 * card_width - 2 * card_gap) / 2

for i, (name, subtitle, price, unit, features_text, is_popular) in enumerate(plans):
    x = start_x + i * (card_width + card_gap)
    
    # Card background
    fill = RGBColor(0, 180, 180) if is_popular else LIGHT_GRAY
    add_rounded_rectangle(slide5, x, 1.5, card_width, 5.3, fill)
    
    # Popular badge
    if is_popular:
        add_rounded_rectangle(slide5, x + 0.8, 1.3, 2.2, 0.5, ACCENT_ORANGE)
        add_text_box(slide5, x + 0.8, 1.35, 2.2, 0.4, "最受歡迎", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    # Plan name
    add_text_box(slide5, x, 1.8, card_width, 0.5, name, font_size=26, bold=True, 
                 color=WHITE if is_popular else PRIMARY_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide5, x, 2.3, card_width, 0.4, subtitle, font_size=16, 
                 color=WHITE if is_popular else DARK_GRAY, align=PP_ALIGN.CENTER)
    
    # Price
    price_color = WHITE if is_popular else SECONDARY_TEAL
    add_text_box(slide5, x, 2.9, card_width, 0.8, price, font_size=36, bold=True, 
                 color=price_color, align=PP_ALIGN.CENTER)
    if unit:
        add_text_box(slide5, x, 3.6, card_width, 0.4, unit, font_size=16, 
                     color=price_color, align=PP_ALIGN.CENTER)
    
    # Features
    add_text_box(slide5, x + 0.3, 4.2, card_width - 0.6, 2.5, features_text, font_size=14, 
                 color=WHITE if is_popular else DARK_GRAY, align=PP_ALIGN.CENTER)

# Free trial note
add_text_box(slide5, 0, 6.9, 13.333, 0.4, "所有方案：14日免費試用，無需信用卡", font_size=16, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# Area2 footer
add_text_box(slide5, 0.3, 7.0, 2, 0.3, "Area2", font_size=12, color=SECONDARY_TEAL)

# ============== SLIDE 6: Who is it for ==============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_rectangle(slide6, 0, 0, 13.333, 1.2, PRIMARY_BLUE)
add_text_box(slide6, 0.5, 0.3, 12, 0.7, "適合邊類人?", font_size=36, bold=True, color=WHITE)

# Target audience cards
audiences = [
    ("中小企老闆", "零售、餐飲、服務行業"),
    ("零售店東主", "多店管理、庫存追蹤"),
    ("餐廳東主", "餐牌更新、訂位管理"),
    ("美容院", "預約提醒、客戶管理"),
    ("地產代理", "放盤文案、客戶跟進"),
    ("會計/物流", "單據追蹤、派送更新"),
]

for i, (title, desc) in enumerate(audiences):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 1.6 + row * 2.6
    
    add_rounded_rectangle(slide6, x, y, 3.8, 2.2, SECONDARY_TEAL)
    add_text_box(slide6, x, y + 0.5, 3.8, 0.6, title, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide6, x, y + 1.2, 3.8, 0.8, desc, font_size=15, color=WHITE, align=PP_ALIGN.CENTER)

# Area2 footer
add_text_box(slide6, 0.3, 7.0, 2, 0.3, "Area2", font_size=12, color=SECONDARY_TEAL)

# ============== SLIDE 7: Why Choose SecrexAI ==============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_rectangle(slide7, 0, 0, 13.333, 1.2, PRIMARY_BLUE)
add_text_box(slide7, 0.5, 0.3, 12, 0.7, "為什麼選擇 SecrexAI?", font_size=36, bold=True, color=WHITE)

# Security badges
badges = [
    ("ISO 27001", "企業級資安認證"),
    ("SOC 2 Type II", "服務安全審計"),
    ("PDPO 合規", "香港私隱條例"),
    ("端對端加密", "帳號完全隔離"),
]

badge_y = 1.5
for i, (title, desc) in enumerate(badges):
    x = 0.8 + i * 3.1
    add_rounded_rectangle(slide7, x, badge_y, 2.8, 1.3, PRIMARY_BLUE)
    add_text_box(slide7, x, badge_y + 0.2, 2.8, 0.5, title, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide7, x, badge_y + 0.7, 2.8, 0.5, desc, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Additional benefits
benefits = [
    "數據存亞太區雲端（Tencent Cloud 企業級數據中心）",
    "30分鐘個人化 Onboarding，調教 AI 認識你的工作習慣",
    "WhatsApp 廣東話技術支援，辦公時間無限次",
    "一個平台，60+ 技能全包，一個訂閱解決所有工具需求",
]

add_bullet_list(slide7, 0.8, 3.2, 11.5, 3.5, benefits, font_size=18, color=DARK_GRAY)

# Area2 footer
add_text_box(slide7, 0.3, 7.0, 2, 0.3, "Area2", font_size=12, color=SECONDARY_TEAL)

# ============== SLIDE 8: Call to Action ==============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

# Background
add_rectangle(slide8, 0, 0, 13.333, 7.5, PRIMARY_BLUE)
add_rectangle(slide8, 0, 5, 13.333, 2.5, SECONDARY_TEAL)

# Main CTA
add_text_box(slide8, 1, 2, 11, 1, "立即開始", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide8, 1, 3.2, 11, 0.8, "14 日免費試用", font_size=36, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# Sub text
add_text_box(slide8, 1, 4.2, 11, 0.6, "無需信用卡，直接 WhatsApp 開始", font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

# Contact info
add_text_box(slide8, 1, 5.5, 11, 0.5, "聯絡 Area2 獲取更多資訊", font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide8, 1, 6.2, 11, 0.4, "WhatsApp: 9565 1459  |  info@area2.com", font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

# Area2 logo text
add_text_box(slide8, 0.5, 7.0, 2, 0.3, "Area2", font_size=14, bold=True, color=WHITE)

# Save presentation
output_dir = "/home/ubuntu/.openclaw/workspace/data/projects/secrexai-ppt"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "SecrexAI_Introduction.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
